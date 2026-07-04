#!/usr/bin/env python3
"""build_deck.py -- Generate a complete PPTX + PDF deck from report.json.

Uses the locked Riso Field-Study design language. Generates filled SVG slides,
rasterizes via headless Chrome at 2x, assembles PPTX (python-pptx) + PDF (Pillow).

Usage:
    python tools/build_deck.py --input runs/2026-07-04/report.json [--out reports/]

Dependencies (deck path only, not core):
    pip install python-pptx Pillow
"""

import argparse
import glob
import os
import shutil
import subprocess
import sys
import tempfile
import textwrap
from html import escape as _esc

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, REPO)
from scout.render_report import (
    load_model_from_json, STAGE_NAMES, LEVEL_NAMES, SCOPE_LABELS,
    ReportModel, Recommendation,
)

# ── Design tokens (Riso Field-Study, locked) ───────────────────────────────
PAPER, INK, BLUE, CORAL = "#F9F7F2", "#1B1F2A", "#2438E0", "#EE5340"
SLATE, GHOST, BORDER, WHITE = "#55596A", "#8A8A93", "#E3E1DC", "#FFFFFF"
TF = "'Arial Black', 'Helvetica Neue', Arial, sans-serif"
BF = "Arial, Helvetica, sans-serif"
MF = "Consolas, 'Courier New', monospace"
W, H = 1280, 720


# ── SVG micro-helpers ───────────────────────────────────────────────────────

def e(t: str) -> str:
    return _esc(str(t))


def _hdr() -> str:
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}">'
        f'<rect width="{W}" height="{H}" fill="{PAPER}"/>'
        '<defs><pattern id="grid" width="20" height="20" patternUnits="userSpaceOnUse">'
        f'<path d="M 20 0 L 0 0 0 20" fill="none" stroke="{INK}" stroke-opacity="0.06" stroke-width="0.5"/>'
        f'</pattern></defs><rect width="{W}" height="{H}" fill="url(#grid)"/>'
    )


def _ftr(page: int, caption: str = "") -> str:
    parts = [f'<line x1="60" y1="670" x2="1220" y2="670" stroke="{BORDER}" stroke-width="1"/>']
    if caption:
        parts.append(f'<text x="60" y="695" font-family="{MF}" font-size="13" fill="{GHOST}">{e(caption)}</text>')
    parts.append(f'<text x="1220" y="695" font-family="{MF}" font-size="13" fill="{GHOST}" text-anchor="end">{page}</text>')
    return "\n".join(parts)


def _t(x, y, txt, font=BF, sz=18, fill=INK, w="400", anc="start", ls=None, op=None):
    a = f'x="{x}" y="{y}" font-family="{font}" font-size="{sz}" font-weight="{w}" fill="{fill}"'
    if anc != "start":
        a += f' text-anchor="{anc}"'
    if ls:
        a += f' letter-spacing="{ls}"'
    if op is not None:
        a += f' fill-opacity="{op}"'
    return f"<text {a}>{e(txt)}</text>"


def _tblock(x, y, lines, font=BF, sz=18, fill=INK, lh=None):
    lh = lh or int(sz * 1.45)
    return "\n".join(_t(x, y + i * lh, ln, font=font, sz=sz, fill=fill) for i, ln in enumerate(lines))


def _wrap(txt, w=65):
    return textwrap.wrap(txt, width=w) or [""]


def _slider(val, color, x, y, cw=30, ch=18):
    return "\n".join(
        f'<rect x="{x + i * cw}" y="{y}" width="{cw}" height="{ch}" '
        f'fill="{color if i < val else "none"}" stroke="{color}" stroke-width="0.5"/>'
        for i in range(5)
    )


def _cites_text(cites, maxlen=90):
    labels = [c.label for c in cites]
    txt = " | ".join(labels)
    return txt[:maxlen] + "..." if len(txt) > maxlen else txt


# ── Slide generators ───────────────────────────────────────────────────────

def slide_cover(m: ReportModel) -> str:
    subj_line = m.subject_line
    if len(subj_line) > 75:
        subj_line = subj_line[:72] + "..."
    return "\n".join([
        _hdr(),
        _t(70, 470, "CHAOS", TF, 200, BLUE, "800", ls="-4", op=0.06),
        _t(470, 600, "METAMORPHOSIS", TF, 120, CORAL, "800", ls="-2", op=0.06),
        f'<line x1="60" y1="60" x2="1220" y2="60" stroke="{BORDER}" stroke-width="1"/>',
        _t(60, 95, f"AI CHAOS SCOUT // {m.subject_name.upper()} // {m.date}", MF, 15, SLATE, ls="3"),
        f'<rect x="60" y="140" width="64" height="10" fill="{BLUE}"/>',
        _t(58, 250, "AI Chaos Scout", TF, 86, INK, "800", ls="-2"),
        _t(60, 300, "Weekly Chaos Briefing", TF, 30, BLUE, "700", ls="1"),
        _t(60, 390, m.subject_name, TF, 44, INK, "800"),
        f'<rect x="60" y="410" width="1160" height="2" fill="{CORAL}"/>',
        _t(60, 452, subj_line, BF, 22, SLATE),
        # Metamorphosis strip
        f'<ellipse cx="280" cy="580" rx="18" ry="24" fill="none" stroke="{BLUE}" stroke-width="2"/>',
        f'<path d="M 430,565 Q 445,555 460,565 Q 475,575 490,565 Q 505,555 520,565" fill="none" stroke="{BLUE}" stroke-width="2"/>',
        f'<ellipse cx="680" cy="575" rx="16" ry="24" fill="none" stroke="{CORAL}" stroke-width="2"/>',
        f'<line x1="680" y1="551" x2="680" y2="541" stroke="{CORAL}" stroke-width="1.5"/>',
        f'<path d="M 858,580 L 838,555 L 828,575" fill="none" stroke="{CORAL}" stroke-width="2"/>',
        f'<path d="M 892,580 L 912,555 L 922,575" fill="none" stroke="{CORAL}" stroke-width="2"/>',
        f'<path d="M 1050,575 L 1030,550 L 1020,575 L 1050,575" fill="none" stroke="{CORAL}" stroke-width="2"/>',
        f'<path d="M 1090,575 L 1110,550 L 1120,575 L 1090,575" fill="none" stroke="{CORAL}" stroke-width="2"/>',
        *[_t(cx, 627, nm, MF, 13, GHOST, anc="middle")
          for cx, nm in [(280, "egg"), (475, "larva"), (680, "chrysalis"), (875, "emergence"), (1070, "imago")]],
        *[_t(cx, 645, f"L{i}", MF, 11, GHOST, anc="middle") for cx, i in [(280, 1), (475, 2), (680, 3), (875, 4), (1070, 5)]],
        *[f'<line x1="{a}" y1="580" x2="{b}" y2="580" stroke="{BORDER}" stroke-width="1"/>'
          for a, b in [(310, 420), (530, 655), (705, 820), (930, 1015)]],
        f'<line x1="60" y1="675" x2="1220" y2="675" stroke="{BORDER}" stroke-width="1"/>',
        _t(1220, 700, m.date, MF, 13, GHOST, anc="end"),
        "</svg>",
    ])


def slide_tldr(m: ReportModel, page: int) -> str:
    parts = [_hdr()]
    parts.append(f'<line x1="60" y1="60" x2="1220" y2="60" stroke="{BORDER}" stroke-width="1"/>')
    parts.append(_t(60, 95, "TL;DR", MF, 15, SLATE, ls="3"))
    parts.append(f'<rect x="60" y="130" width="64" height="10" fill="{BLUE}"/>')
    parts.append(_t(60, 195, "Executive Summary", TF, 42, INK, "800"))
    y = 260
    for i, bullet in enumerate(m.tldr):
        lines = _wrap(bullet, 80)
        parts.append(f'<rect x="60" y="{y - 5}" width="8" height="8" fill="{BLUE if i == 0 else CORAL if i == 2 else INK}"/>')
        parts.append(_tblock(85, y, lines, sz=18, fill=INK, lh=26))
        y += len(lines) * 26 + 30
    parts.append(_ftr(page, f"AI Chaos Scout \u2014 {m.subject_name} \u2014 {m.date}"))
    parts.append("</svg>")
    return "\n".join(parts)


def slide_scoreboard(m: ReportModel, page: int) -> str:
    parts = [_hdr()]
    parts.append(f'<line x1="60" y1="60" x2="1220" y2="60" stroke="{BORDER}" stroke-width="1"/>')
    parts.append(_t(60, 95, "SCOREBOARD", MF, 15, SLATE, ls="3"))
    parts.append(f'<rect x="60" y="130" width="64" height="10" fill="{BLUE}"/>')
    parts.append(_t(60, 185, "Recommendations at a Glance", TF, 36, INK, "800"))
    # Table header
    hy = 220
    cols = [("Code", 60), ("Title", 130), ("Feas", 660), ("Evid", 760), ("Imp", 860), ("Disr", 960), ("Stage", 1060)]
    parts.append(f'<line x1="60" y1="{hy + 5}" x2="1220" y2="{hy + 5}" stroke="{INK}" stroke-width="1"/>')
    for label, cx in cols:
        parts.append(_t(cx, hy, label, MF, 12, SLATE, w="400"))
    ry = hy + 30
    for rec in m.recommendations:
        color = CORAL if rec.track == "B" else BLUE
        title = rec.title[:45] + "..." if len(rec.title) > 48 else rec.title
        parts.append(_t(60, ry, rec.code, MF, 14, color, w="700"))
        parts.append(_t(130, ry, title, BF, 14, INK))
        parts.append(_slider(rec.feasibility, color, 660, ry - 13, 18, 14))
        parts.append(_slider(rec.evidence, color, 760, ry - 13, 18, 14))
        parts.append(_slider(rec.impact, color, 860, ry - 13, 18, 14))
        parts.append(_slider(rec.disruption, color, 960, ry - 13, 18, 14))
        stage = STAGE_NAMES.get(rec.disruption, "")
        parts.append(_t(1060, ry, stage, MF, 13, GHOST))
        parts.append(f'<line x1="60" y1="{ry + 12}" x2="1220" y2="{ry + 12}" stroke="{BORDER}" stroke-width="0.5"/>')
        ry += 42
    # Legend
    ly = ry + 20
    parts.append(f'<rect x="60" y="{ly}" width="12" height="12" fill="{BLUE}"/>')
    parts.append(_t(80, ly + 11, "Track A \u2014 Realistic (disruption 1\u20132)", MF, 12, SLATE))
    parts.append(f'<rect x="400" y="{ly}" width="12" height="12" fill="{CORAL}"/>')
    parts.append(_t(420, ly + 11, "Track B \u2014 Metamorphosis (disruption 3\u20135)", MF, 12, SLATE))
    parts.append(_ftr(page))
    parts.append("</svg>")
    return "\n".join(parts)


def slide_ladder(page: int) -> str:
    parts = [_hdr()]
    parts.append(f'<line x1="60" y1="60" x2="1220" y2="60" stroke="{BORDER}" stroke-width="1"/>')
    parts.append(_t(60, 95, "THE CHAOS METER", MF, 15, SLATE, ls="3"))
    parts.append(f'<rect x="60" y="130" width="64" height="10" fill="{BLUE}"/>')
    parts.append(_t(60, 185, "Disruption Levels 1\u20135 as Metamorphosis Stages", TF, 32, INK, "800"))
    parts.append(_t(60, 220, "the realism ladder \u2014 scope of the bet increases left to right", BF, 18, SLATE))

    stages = [
        (1, "egg", "Incremental", "tweak an existing surface", BLUE),
        (2, "larva", "Adjacent", "feature-level bet", BLUE),
        (3, "chrysalis", "New line", "product/feature bet \u2014 reversible", CORAL),
        (4, "emergence", "Model shift", "line-of-business / model bet", CORAL),
        (5, "imago", "Metamorphosis", "bet-the-company", CORAL),
    ]
    cw = 210
    gap = 20
    x0 = 75
    for i, (lvl, stage, name, scope, color) in enumerate(stages):
        cx = x0 + i * (cw + gap)
        # Card
        parts.append(f'<rect x="{cx}" y="270" width="{cw}" height="340" fill="{WHITE}" stroke="{color}" stroke-width="1"/>')
        parts.append(f'<rect x="{cx}" y="270" width="{cw}" height="4" fill="{color}"/>')
        # Level number
        parts.append(_t(cx + cw // 2, 330, f"L{lvl}", TF, 48, color, "800", anc="middle"))
        # Stage name
        parts.append(_t(cx + cw // 2, 370, stage.upper(), MF, 16, color, anc="middle"))
        # Meaning
        parts.append(_t(cx + cw // 2, 415, name, TF, 22, INK, "700", anc="middle"))
        # Scope
        scope_lines = _wrap(scope, 22)
        parts.append(_tblock(cx + 15, 460, scope_lines, MF, 13, SLATE, lh=18))
        # Track label
        track = "Track A" if lvl <= 2 else "Track B"
        parts.append(_t(cx + cw // 2, 590, track, MF, 12, color, anc="middle"))
        # Arrow to next
        if i < 4:
            ax = cx + cw
            parts.append(f'<line x1="{ax + 2}" y1="440" x2="{ax + gap - 2}" y2="440" stroke="{BORDER}" stroke-width="1"/>')
            parts.append(f'<path d="M {ax + gap - 8},436 L {ax + gap - 2},440 L {ax + gap - 8},444" fill="none" stroke="{BORDER}" stroke-width="1"/>')

    parts.append(_ftr(page, "disruption 1\u20132 = Track A (realistic) \u00b7 disruption 3\u20135 = Track B (metamorphosis)"))
    parts.append("</svg>")
    return "\n".join(parts)


def slide_rec(rec: Recommendation, page: int, fig: int) -> str:
    is_b = rec.track == "B"
    color = CORAL if is_b else BLUE
    track_label = "TRACK B \u2014 METAMORPHOSIS" if is_b else "TRACK A \u2014 REALISTIC"
    stage = STAGE_NAMES.get(rec.disruption, "")
    scope = SCOPE_LABELS.get(rec.disruption, "")

    parts = [_hdr()]
    # Track header
    parts.append(_t(60, 55, track_label, MF, 14, color, ls="2"))
    parts.append(f'<line x1="60" y1="65" x2="1220" y2="65" stroke="{color}" stroke-width="0.5"/>')
    # Card
    parts.append(f'<rect x="60" y="85" width="1160" height="530" fill="{WHITE}" stroke="{color}" stroke-width="0.5"/>')
    if is_b:
        parts.append(f'<rect x="60" y="85" width="4" height="530" fill="{color}"/>')
    # Header
    title = rec.title[:50] + "..." if len(rec.title) > 53 else rec.title
    parts.append(_t(90, 128, rec.code, MF, 22, color, w="700"))
    parts.append(_t(160, 128, title, TF, 28, color, w="700"))
    level_label = LEVEL_NAMES.get(rec.disruption, "")
    badge = f"Level {rec.disruption} \u2014 {stage}"
    if is_b:
        badge += f" ({scope})"
    parts.append(_t(1190, 128, badge, MF, 13, color, anc="end"))
    parts.append(f'<line x1="90" y1="145" x2="1190" y2="145" stroke="{color}" stroke-width="0.5"/>')

    # Body
    body_lines = _wrap(rec.body, 75)[:4]
    parts.append(_tblock(90, 178, body_lines, BF, 17, INK, lh=24))
    body_end = 178 + len(body_lines) * 24 + 8

    # Sliders
    sy = body_end + 5
    parts.append(f'<line x1="90" y1="{sy}" x2="1190" y2="{sy}" stroke="{color}" stroke-width="0.5"/>')
    sy += 25
    sliders = [("Feasibility", rec.feasibility, 90), ("Evidence", rec.evidence, 370),
               ("Impact", rec.impact, 650), ("Disruption", rec.disruption, 930)]
    for label, val, sx in sliders:
        parts.append(_t(sx, sy, label, MF, 13, SLATE))
        parts.append(_slider(val, color, sx + 100, sy - 14, 26, 16))
        parts.append(_t(sx + 235, sy, f"{val}/5", MF, 13, GHOST))
    sy += 25
    parts.append(f'<line x1="90" y1="{sy}" x2="1190" y2="{sy}" stroke="{color}" stroke-width="0.5"/>')

    # Why now
    wny = sy + 30
    parts.append(_t(90, wny, "Why now", TF, 17, color, w="700"))
    wn_text = _cites_text(rec.why_now, 100)
    wn_lines = _wrap(wn_text, 75)[:2]
    parts.append(_tblock(90, wny + 22, wn_lines, BF, 15, INK, lh=20))
    wny += 22 + len(wn_lines) * 20 + 10

    # Why us
    parts.append(_t(90, wny, "Why us", TF, 17, color, w="700"))
    wu_text = _cites_text(rec.why_us, 100)
    wu_lines = _wrap(wu_text, 75)[:2]
    parts.append(_tblock(90, wny + 22, wu_lines, BF, 15, INK, lh=20))
    wny += 22 + len(wu_lines) * 20 + 15

    # Track-specific section
    parts.append(f'<line x1="90" y1="{wny}" x2="1190" y2="{wny}" stroke="{color}" stroke-width="0.5"/>')
    if is_b:
        parts.append(_t(90, wny + 25, "Metamorphosis narrative", TF, 16, color, w="700"))
        meta_lines = _wrap(rec.metamorphosis, 80)[:3]
        parts.append(_tblock(90, wny + 47, meta_lines, BF, 14, INK, lh=19))
        ky = wny + 47 + len(meta_lines) * 19 + 10
        parts.append(f'<rect x="90" y="{ky}" width="1100" height="35" fill="{color}" fill-opacity="0.06" stroke="{color}" stroke-width="0.5"/>')
        parts.append(_t(105, ky + 23, "KILL CRITERIA:", MF, 13, color))
        kc = rec.kill_criteria[:95] + "..." if len(rec.kill_criteria) > 98 else rec.kill_criteria
        parts.append(_t(260, ky + 23, kc, BF, 13, INK))
    else:
        parts.append(_t(90, wny + 25, "First testable step", TF, 16, color, w="700"))
        fs_lines = _wrap(rec.first_step, 80)[:3]
        parts.append(_tblock(90, wny + 47, fs_lines, BF, 14, INK, lh=19))

    cap = f"Fig. {fig} \u2014 {'Metamorphosis' if is_b else 'Specimen'} plate {rec.code}"
    parts.append(_ftr(page, cap))
    parts.append("</svg>")
    return "\n".join(parts)


def slide_themes(themes, start_idx: int, page: int) -> str:
    parts = [_hdr()]
    parts.append(f'<line x1="60" y1="60" x2="1220" y2="60" stroke="{BORDER}" stroke-width="1"/>')
    label = "DEEP DIVE" if start_idx == 0 else "DEEP DIVE (CONT.)"
    parts.append(_t(60, 95, label, MF, 15, SLATE, ls="3"))
    parts.append(f'<rect x="60" y="130" width="64" height="10" fill="{BLUE}"/>')
    parts.append(_t(60, 185, "Digest Themes", TF, 36, INK, "800"))

    y = 230
    for i, theme in enumerate(themes):
        parts.append(f'<rect x="60" y="{y}" width="1160" height="2" fill="{BLUE if (start_idx + i) % 2 == 0 else CORAL}"/>')
        parts.append(_t(60, y + 28, theme.heading, TF, 20, INK, w="700"))
        body_lines = _wrap(theme.body, 85)[:3]
        parts.append(_tblock(60, y + 52, body_lines, BF, 15, INK, lh=21))
        cite_y = y + 52 + len(body_lines) * 21 + 5
        cite_text = " | ".join(c.label[:40] for c in theme.citations[:3])
        parts.append(_t(60, cite_y, cite_text[:110], MF, 11, GHOST))
        y = cite_y + 30

    parts.append(_ftr(page))
    parts.append("</svg>")
    return "\n".join(parts)


def slide_provenance(m: ReportModel, page: int) -> str:
    prov = m.provenance
    parts = [_hdr()]
    parts.append(f'<line x1="60" y1="60" x2="1220" y2="60" stroke="{BORDER}" stroke-width="1"/>')
    parts.append(_t(60, 95, "PROVENANCE", MF, 15, SLATE, ls="3"))
    parts.append(f'<rect x="60" y="130" width="64" height="10" fill="{BLUE}"/>')
    parts.append(_t(60, 185, "Source Inventory & Grounding Note", TF, 32, INK, "800"))

    # Stats card
    parts.append(f'<rect x="60" y="230" width="1160" height="120" fill="{WHITE}" stroke="{BORDER}" stroke-width="1"/>')
    stats = [
        ("SCANNED", str(prov.scanned), BLUE),
        ("SOURCES", str(prov.sources), BLUE),
        ("RELEVANT", str(prov.relevant), CORAL),
        ("SOFT-FAILED", str(len(prov.soft_failed)) if prov.soft_failed else "0", GHOST),
    ]
    for j, (label, val, color) in enumerate(stats):
        sx = 120 + j * 280
        parts.append(_t(sx, 280, val, TF, 48, color, "800", anc="middle"))
        parts.append(_t(sx, 310, label, MF, 13, SLATE, anc="middle"))

    if prov.soft_failed:
        parts.append(_t(60, 385, "Soft-failed sources:", BF, 16, CORAL))
        parts.append(_t(250, 385, ", ".join(prov.soft_failed), MF, 15, INK))
    else:
        parts.append(_t(60, 385, "All sources responded successfully.", BF, 16, INK))

    # Grounding note
    parts.append(f'<rect x="60" y="420" width="1160" height="180" fill="{BLUE}" fill-opacity="0.04" stroke="{BLUE}" stroke-width="0.5"/>')
    parts.append(_t(90, 455, "GROUNDING RULE", MF, 14, BLUE, w="700"))
    note_lines = [
        "Every recommendation cites at least one 'why now' (a scouted item from this run)",
        "AND at least one 'why us' (a compiled-truth fact about the subject company).",
        "Recommendations failing this gate are deleted before rendering.",
        "Transformation ideas must emerge from a real trend x a real fact.",
    ]
    parts.append(_tblock(90, 485, note_lines, BF, 15, INK, lh=24))

    parts.append(_ftr(page, f"AI Chaos Scout \u2014 {m.subject_name} \u2014 {m.date}"))
    parts.append("</svg>")
    return "\n".join(parts)


# ── Rasterization + assembly ───────────────────────────────────────────────

def _find_chrome() -> str:
    candidates = [
        "/opt/.devin/chrome/chrome/linux-133.0.6943.126/chrome-linux64/chrome",
    ]
    # Also search for any chrome binary under /opt
    candidates += sorted(glob.glob("/opt/**/chrome", recursive=True))
    for c in candidates:
        if os.path.isfile(c) and os.access(c, os.X_OK):
            return c
    raise RuntimeError("No Chrome binary found for headless rendering")


def rasterize_svgs(svg_strings: list[str], out_dir: str, chrome: str) -> list[str]:
    """Write SVGs to temp files, screenshot each via headless Chrome, return PNG paths."""
    png_paths = []
    with tempfile.TemporaryDirectory() as tmp:
        for i, svg in enumerate(svg_strings):
            svg_path = os.path.join(tmp, f"slide_{i:03d}.svg")
            png_path = os.path.join(out_dir, f"slide_{i:03d}.png")
            with open(svg_path, "w", encoding="utf-8") as f:
                f.write(svg)
            profile = os.path.join(tmp, f"profile_{i}")
            cmd = [
                chrome, "--headless=new", "--no-sandbox", "--disable-gpu",
                "--disable-dev-shm-usage", f"--user-data-dir={profile}",
                "--hide-scrollbars", "--force-device-scale-factor=2",
                f"--window-size={W},{H}", f"--screenshot={png_path}",
                f"file://{svg_path}",
            ]
            subprocess.run(cmd, capture_output=True, timeout=60)
            if os.path.isfile(png_path):
                png_paths.append(png_path)
            else:
                print(f"  WARNING: slide {i} failed to rasterize", file=sys.stderr)
    return png_paths


def assemble_pdf(png_paths: list[str], pdf_path: str):
    """Combine PNGs into a multi-page PDF via Pillow."""
    from PIL import Image
    images = []
    for p in png_paths:
        img = Image.open(p).convert("RGB")
        images.append(img)
    if not images:
        raise RuntimeError("No slides to assemble")
    images[0].save(pdf_path, save_all=True, append_images=images[1:], resolution=150)
    print(f"  PDF: {pdf_path} ({len(images)} pages)")


def assemble_pptx(png_paths: list[str], pptx_path: str):
    """Build a PPTX with one full-bleed image per slide."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank
    for png in png_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(png, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    prs.save(pptx_path)
    print(f"  PPTX: {pptx_path} ({len(png_paths)} slides)")


# ── Main ───────────────────────────────────────────────────────────────────

def build_deck(input_path: str, out_dir: str | None = None) -> dict[str, str]:
    model = load_model_from_json(input_path)
    out_dir = out_dir or os.path.join(REPO, "reports")
    os.makedirs(out_dir, exist_ok=True)

    print(f"Building deck for {model.subject_name} ({model.date})...")
    chrome = _find_chrome()
    print(f"  Chrome: {chrome}")

    # Generate slides
    track_a = [r for r in model.recommendations if r.track == "A"]
    track_b = [r for r in model.recommendations if r.track == "B"]
    themes = model.digest_themes

    slides = []
    page = 1

    # 1. Cover
    slides.append(slide_cover(model))
    page += 1

    # 2. TL;DR
    slides.append(slide_tldr(model, page))
    page += 1

    # 3. Scoreboard
    slides.append(slide_scoreboard(model, page))
    page += 1

    # 4. Disruption ladder
    slides.append(slide_ladder(page))
    page += 1

    # 5-N. Track A rec cards
    fig = 1
    for rec in track_a:
        slides.append(slide_rec(rec, page, fig))
        page += 1
        fig += 1

    # N+1-M. Track B rec cards
    for rec in track_b:
        slides.append(slide_rec(rec, page, fig))
        page += 1
        fig += 1

    # M+1. Digest themes (split across 1-2 slides)
    if themes:
        mid = min(3, len(themes))
        slides.append(slide_themes(themes[:mid], 0, page))
        page += 1
        if len(themes) > mid:
            slides.append(slide_themes(themes[mid:], mid, page))
            page += 1

    # Last. Provenance
    slides.append(slide_provenance(model, page))

    print(f"  {len(slides)} slides generated")

    # Rasterize
    with tempfile.TemporaryDirectory() as png_dir:
        png_paths = rasterize_svgs(slides, png_dir, chrome)
        print(f"  {len(png_paths)} slides rasterized")

        if not png_paths:
            raise RuntimeError("No slides rasterized successfully")

        # Assemble
        date = model.date
        subj = model.subject_name.lower().replace(" ", "-")
        pdf_path = os.path.join(out_dir, f"chaos-deck-{subj}-{date}.pdf")
        pptx_path = os.path.join(out_dir, f"chaos-deck-{subj}-{date}.pptx")

        assemble_pdf(png_paths, pdf_path)
        assemble_pptx(png_paths, pptx_path)

    return {"pdf": pdf_path, "pptx": pptx_path}


def main():
    ap = argparse.ArgumentParser(description="Build a completed PPTX + PDF deck from report.json")
    ap.add_argument("--input", required=True, help="Path to report.json")
    ap.add_argument("--out", default=None, help="Output directory (default: reports/)")
    args = ap.parse_args()
    paths = build_deck(args.input, args.out)
    for kind, p in paths.items():
        print(f"{kind}: {p}")


if __name__ == "__main__":
    main()
